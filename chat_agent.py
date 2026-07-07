import json
import base64
import io
import re
from dataclasses import dataclass
from typing import List, Dict, Any, Optional, AsyncGenerator, Union
from openai import AzureOpenAI

# === 防禦性引入三大文件解析套件 ===
try:
    import docx
except ImportError:
    docx = None

try:
    import pdfplumber
except ImportError:
    pdfplumber = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None


@dataclass
class FileInput:
    """標準化檔案輸入結構"""
    filename: str
    file_base64: str


class ContextManager:
    """處理對話歷史紀錄的邏輯"""
    
    @staticmethod
    def count_tokens(text: str) -> int:
        return len(text.split())

    @staticmethod
    def get_last_n_pairs_strict(history: List[Dict], n: int) -> List[Dict]:
        pairs = []
        i = 0
        L = len(history)
        while i < L-1:
            if history[i]['role'] == 'user' and history[i+1]['role'] == 'assistant':
                pairs.append([history[i], history[i+1]])
                i += 2
            else:
                i += 1
        output = []
        for pair in pairs[-n:]:
            output.extend(pair)
        if L > 0:
            last_is_user = history[-1]['role'] == 'user'
            if last_is_user:
                seen = any(msg == history[-1] for msg in output[::-1])
                if not seen:
                    output.append(history[-1])
        return output

    @staticmethod
    def get_last_topic_chunk(history: List[Dict], max_rounds: int = 5) -> List[Dict]:
        chunk = []
        rounds = 0
        i = len(history) - 1
        while i >= 0 and history[i]['role'] != 'user':
            i -= 1
        while i >= 0 and rounds < max_rounds:
            current = []
            if history[i]['role'] == 'user':
                current.insert(0, history[i])
                i -= 1
                if i >= 0 and history[i]['role'] == 'assistant':
                    current.insert(0, history[i])
                    i -= 1
            else:
                i -= 1
                continue
            chunk.insert(0, current)
            rounds += 1
        return [msg for pair in chunk for msg in pair]

    @classmethod
    def get_ai_context_auto(cls, history: List[Dict], n: int = 5, token_threshold: int = 3000, short_reply_limit: int = 30) -> List[Dict]:
        if not history:
            return []
        last_roles = [msg['role'] for msg in history[-4:]]
        last_msg = history[-1]

        if last_msg['role'] == 'user' or last_roles[-2:] == ['user', 'user']:
            return cls.get_last_topic_chunk(history, max_rounds=n)

        for msg in reversed(history):
            if msg['role'] == 'user' and cls.count_tokens(str(msg.get('content', ''))) > token_threshold:
                return cls.get_last_topic_chunk(history, max_rounds=n)
            break

        for msg in reversed(history):
            if msg['role'] == 'assistant' and cls.count_tokens(str(msg.get('content', ''))) < short_reply_limit:
                return cls.get_last_topic_chunk(history, max_rounds=n)
            break

        return cls.get_last_n_pairs_strict(history, n=n)


class DocumentParser:
    """處理檔案解析邏輯 (純記憶體不落地)"""
    
    @staticmethod
    def parse(file_obj: FileInput) -> Dict[str, Any]:
        """
        回傳字典：
        {
            "type": "text" 或 "multimodal",
            "content": 解析出來的純文字 或 包含圖片的陣列
        }
        """
        result = {"type": "text", "content": ""}
        pdf_contents = []

        try:
            file_bytes = base64.b64decode(file_obj.file_base64)
            file_stream = io.BytesIO(file_bytes)
            filename = file_obj.filename.lower()

            if filename.endswith(".docx"):
                if docx:
                    from docx.oxml.table import CT_Tbl
                    from docx.oxml.text.paragraph import CT_P
                    from docx.table import Table
                    from docx.text.paragraph import Paragraph

                    doc = docx.Document(file_stream)
                    content_lines = []
                    for child in doc.element.body:
                        if isinstance(child, CT_P):
                            para = Paragraph(child, doc)
                            if para.text.strip():
                                content_lines.append(para.text.strip())
                        elif isinstance(child, CT_Tbl):
                            table = Table(child, doc)
                            content_lines.append("\n--- [表格資料開始] ---")
                            for row in table.rows:
                                row_data = [cell.text.replace('\n', ' ').strip() for cell in row.cells]
                                content_lines.append(" | ".join(row_data))
                            content_lines.append("--- [表格資料結束] ---\n")
                    
                    extracted_text = "\n".join(content_lines)
                    result["content"] = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '', extracted_text)
                else:
                    result["content"] = "[系統提示：環境尚未安裝 python-docx，無法解析 Word]"

            elif filename.endswith(".pdf"):
                if pdfplumber:
                    with pdfplumber.open(file_stream) as pdf:
                        image_count = 0
                        max_images = 50  # 限制最多傳送 50 張圖片
                        
                        for page_num, page in enumerate(pdf.pages, 1):
                            page_text = page.extract_text()
                            if page_text:
                                pdf_contents.append({
                                    "type": "text",
                                    "text": f"--- Page {page_num} 文字 ---\n{page_text}"
                                })
                            
                            # 控制圖片數量不超過 API 上限
                            if image_count < max_images:
                                img = page.to_image(resolution=100).original
                                buffered = io.BytesIO()
                                img.save(buffered, format="JPEG")
                                img_base64 = base64.b64encode(buffered.getvalue()).decode('utf-8')
                                pdf_contents.append({
                                    "type": "image_url",
                                    "image_url": {"url": f"data:image/jpeg;base64,{img_base64}"}
                                })
                                image_count += 1
                            elif image_count == max_images:
                                # 剛好達到上限時，加入一次系統提示
                                pdf_contents.append({
                                    "type": "text",
                                    "text": "\n[系統提示：因 OpenAI API 限制，超過 50 頁的圖檔已省略，後續頁面僅保留純文字解析]"
                                })
                                image_count += 1 # 確保提示只加入一次
                                
                    result["type"] = "multimodal"
                    result["content"] = pdf_contents
                else:
                    result["content"] = "[系統提示：環境尚未安裝 pdfplumber，無法解析 PDF]"

            elif filename.endswith((".pptx", ".ppt")):
                if Presentation:
                    pptx_text = []
                    prs = Presentation(file_stream)
                    for slide_num, slide in enumerate(prs.slides, 1):
                        slide_content = []
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                slide_content.append(shape.text.strip())
                        if slide_content:
                            pptx_text.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_content))
                    result["content"] = "\n".join(pptx_text)
                else:
                    result["content"] = "[系統提示：環境尚未安裝 python-pptx，無法解析 PPTX]"
            else:
                result["content"] = f"[系統提示：不支援的檔案格式 {filename}]"
                
        except Exception as e:
            result["content"] = f"[系統提示：檔案解析失敗 ({file_obj.filename})，原因：{str(e)}]"
            
        return result


class AzureChatAgent:
    """整合 Azure OpenAI 的主要對話類別"""
    
    def __init__(self, endpoint: str, api_key: str, api_version: str, deployment_name: str):
        self.deployment = deployment_name
        self.client = AzureOpenAI(
            api_version=api_version,
            azure_endpoint=endpoint,
            api_key=api_key,
        )

    def _inject_file_context(self, history_context: List[Dict], file_obj: FileInput, parsed_data: Dict[str, Any]) -> None:
        """將解析後的檔案內容注入到對話歷史中"""
        file_type = parsed_data["type"]
        content = parsed_data["content"]

        if file_type == "multimodal" and content:
            injected = False
            header_block = {"type": "text", "text": f"\n\n--- 系統已自動解析使用者上傳的附件檔案 ({file_obj.filename}) ---"}
            footer_block = {"type": "text", "text": "\n--- 附件解析內容結束 ---"}
            
            for msg in reversed(history_context):
                if msg.get("role") == "user":
                    original_content = msg.get("content", "")
                    if isinstance(original_content, str):
                        new_content = [{"type": "text", "text": original_content}]
                    elif isinstance(original_content, list):
                        new_content = original_content
                    else:
                        new_content = [{"type": "text", "text": str(original_content)}]
                    
                    new_content.append(header_block)
                    new_content.extend(content)
                    new_content.append(footer_block)
                    
                    msg["content"] = new_content
                    injected = True
                    break
            
            if not injected:
                new_content = [{"type": "text", "text": f"這是使用者剛剛上傳的附件檔案 ({file_obj.filename})，請參考內容："}]
                new_content.extend(content)
                history_context.append({"role": "user", "content": new_content})

        elif file_type == "text" and str(content).strip():
            injected = False
            for msg in reversed(history_context):
                if msg.get("role") == "user":
                    msg["content"] += (
                        f"\n\n--- 系統已自動解析使用者上傳的附件檔案 ({file_obj.filename}) ---\n"
                        f"{content}\n"
                        f"--- 附件解析內容結束 ---"
                    )
                    injected = True
                    break
            
            if not injected:
                history_context.append({
                    "role": "user",
                    "content": f"這是使用者剛剛上傳的附件檔案 ({file_obj.filename})，請參考內容：\n{content}"
                })

    def build_messages(self, prompt: Union[str, List[Dict]], file_obj: Optional[FileInput] = None, sys_prompt: str = "") -> List[Dict]:
        """
        僅負責建構、裁切歷史紀錄與注入檔案，回傳準備好的 messages 陣列。
        不直接呼叫 API，讓外部 (main.py/routers.py) 保有最大彈性與控制權。
        """
        # 1. 處理歷史紀錄
        history = None
        if isinstance(prompt, str):
            try:
                history = json.loads(prompt)
            except Exception:
                pass
        else:
            history = prompt

        history_context = []
        if isinstance(history, list):
            h = ContextManager.get_ai_context_auto(history, n=5)
            if isinstance(h, list):
                history_context = h

        # 2. 解析並注入檔案
        if file_obj and history_context:
            parsed_data = DocumentParser.parse(file_obj)
            self._inject_file_context(history_context, file_obj, parsed_data)

        # 3. 準備發送給 LLM 的訊息
        messages = []
        if sys_prompt:
            messages.append({"role": "system", "content": sys_prompt})
        messages.extend(history_context)

        return messages
